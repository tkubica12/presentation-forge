"""Adapter: presentation-forge model -> hve-core content tree.

The vendored `microsoft/hve-core` PowerPoint skill at
`skills/pptx-render/` consumes a folder layout shaped like:

    content/
      global/style.yaml
      slide-001/
        content.yaml
        images/<png files used by this slide>
      slide-002/
        ...

Our authoring layer is `slides.md` + `theme.yaml` + `selections.json`
(plus generated PNGs under `<folder>/build/images/<image_ref>/<model>/...`).

`materialize_workspace` writes the upstream-shaped content tree into a
build-scratch directory and returns the paths needed to invoke
`build_deck.py`. Two flavours are supported:

  * ``draft``  — one slide per existing image variant (with a label
                 suffix) so the user can pick a winner, plus textual
                 slides verbatim
  * ``final``  — one slide per spec slide, using the user's selection
                 (or a placeholder note for slides whose image is
                 missing or unselected)

Each slide is rendered onto a named template layout via the mapping in
`theme.yaml.layouts`; placeholders are filled where the layout has them,
and remaining content (images, decorative elements) is emitted as
hve-core ``elements:``.
"""
from __future__ import annotations

import re
import shutil
from dataclasses import dataclass
from pathlib import Path
from typing import Any

import yaml

from .layouts import Layout
from .slides_parser import Slide
from .spec import Presentation, Selection, Theme


# Default logical-name -> Microsoft Azure template layout-name mapping.
# Used as a fallback when theme.yaml does not provide its own `layouts:`
# block. Picks layouts that exist in
# .templates/Microsoft-Azure-PowerPoint-Template-2604.potx and play the
# role expected by each Layout enum value.
DEFAULT_AZURE_LAYOUTS: dict[str, str] = {
    Layout.TITLE.value: "Title Slide 1",
    # Cover reuses the full-bleed photo layout; we replace its content with
    # an image + dark left panel + title/subtitle textboxes (placeholders
    # are intentionally not used).
    Layout.COVER.value: "Photo full bleed lower title",
    Layout.SECTION_DIVIDER.value: "Section Slide 1",
    Layout.BULLETS.value: "Title and Content",
    Layout.BULLETS_WITH_IMAGE.value: "Photo Slide 1",
    Layout.FULL_BLEED_IMAGE.value: "Photo full bleed lower title",
    Layout.TWO_COLUMN.value: "Two Column Bullet text",
    Layout.QUOTE.value: "Quote",
    Layout.COMPARISON.value: "Two Column Bullet text",
    Layout.IMAGE_GRID.value: "Three Filmstrip Photos",
    Layout.IMAGE_SINGLE.value: "Title Only",
    Layout.IMAGE_DUO.value: "Two picture content",
    Layout.APPENDIX_REFERENCES.value: "Title & Non-bulleted text",
}


# Per-layout placeholder index conventions for the Azure template.
# Discovered empirically by inspecting `.templates/...potx`. The fields
# describe which placeholder index inside the layout receives which kind
# of content. None means "no placeholder of that role on this layout".
@dataclass(frozen=True)
class PlaceholderRoles:
    title: int | None = 0
    body: int | None = None       # primary text body / left column
    secondary: int | None = None  # right column or attribution
    picture: int | None = None    # PICTURE placeholder, if any


DEFAULT_PLACEHOLDER_ROLES: dict[str, PlaceholderRoles] = {
    Layout.TITLE.value: PlaceholderRoles(title=0, body=12),
    # Cover: title/subtitle drawn as textboxes inside left panel; no
    # placeholders consumed.
    Layout.COVER.value: PlaceholderRoles(title=None),
    Layout.SECTION_DIVIDER.value: PlaceholderRoles(title=0),
    Layout.BULLETS.value: PlaceholderRoles(title=0, body=10),
    # "Photo Slide 1": title=0 (top-left), body=12 (left below title), picture=10 (right half)
    Layout.BULLETS_WITH_IMAGE.value: PlaceholderRoles(title=0, body=12, picture=10),
    Layout.FULL_BLEED_IMAGE.value: PlaceholderRoles(title=0, picture=10),
    Layout.TWO_COLUMN.value: PlaceholderRoles(title=0, body=12, secondary=13),
    # Quote: we draw body + attribution as explicit textboxes for control
    # over vertical position; placeholders are intentionally bypassed.
    Layout.QUOTE.value: PlaceholderRoles(title=None),
    Layout.COMPARISON.value: PlaceholderRoles(title=0, body=12, secondary=13),
    Layout.IMAGE_GRID.value: PlaceholderRoles(title=0),
    Layout.IMAGE_SINGLE.value: PlaceholderRoles(title=0),
    Layout.IMAGE_DUO.value: PlaceholderRoles(title=0),
    Layout.APPENDIX_REFERENCES.value: PlaceholderRoles(title=0, body=10),
}


# Slide canvas constants (16:9 EMU dimensions of the Azure template).
SLIDE_W = 13.333
SLIDE_H = 7.5


def _placeholder_value(text: str | None) -> str | None:
    """Normalize text destined for a placeholder. Empty -> None."""
    if text is None:
        return None
    s = text.strip()
    return s or None


def _bullets_to_lines(bullets: list[str]) -> list[str] | None:
    items = [b.strip() for b in (bullets or []) if b and b.strip()]
    return items or None


def _resolve_selected_image(pres: Presentation, slide: Slide) -> Path | None:
    if not slide.image_ref:
        return None
    sel = pres.selections.get(slide.slide_id)
    if sel is None:
        return None
    return (
        pres.images_dir / slide.image_ref / sel.model.lower()
        / sel.filename(slide.image_ref)
    )


_VARIANT_RE = re.compile(r"^(?P<ref>.+)_v(?P<v>\d+)_i(?P<i>\d+)\.png$")


def _list_variants(
    pres: Presentation, slide: Slide
) -> list[tuple[str, int, int, Path]]:
    out: list[tuple[str, int, int, Path]] = []
    if not slide.image_ref:
        return out
    base = pres.images_dir / slide.image_ref
    if not base.exists():
        return out
    for model_dir in sorted(p for p in base.iterdir() if p.is_dir()):
        for png in sorted(model_dir.iterdir()):
            m = _VARIANT_RE.match(png.name)
            if m and m.group("ref") == slide.image_ref:
                out.append(
                    (model_dir.name, int(m.group("v")), int(m.group("i")), png)
                )
    return out


def _picture_placeholder_dims(
    template_path: Path | None, layout_name: str, idx: int
) -> tuple[float, float, float, float] | None:
    """Look up the inch coordinates of a picture placeholder in *template_path*.

    Returns (left, top, width, height) in inches, or ``None`` if the layout
    or placeholder cannot be found. We use this to draw an explicit ``image``
    element where the layout would normally place a picture, since hve-core's
    placeholder code does not insert pictures into picture-placeholders.
    """
    if template_path is None or not template_path.exists():
        return None
    # Local import keeps `python-pptx` out of the import chain when callers
    # only need the textual transformation.
    from pptx import Presentation as PptxPresentation
    from pptx.util import Emu

    prs = PptxPresentation(str(template_path))
    for layout in prs.slide_layouts:
        if layout.name != layout_name:
            continue
        for ph in layout.placeholders:
            if ph.placeholder_format.idx == idx:
                return (
                    Emu(ph.left).inches,
                    Emu(ph.top).inches,
                    Emu(ph.width).inches,
                    Emu(ph.height).inches,
                )
    return None


def _compute_fill_crop(
    image_path: Path, target_w: float, target_h: float
) -> dict[str, int] | None:
    """Compute OOXML ``a:srcRect`` crop values for fill-mode display.

    Compares the actual pixel aspect ratio of *image_path* against the
    *target_w*/*target_h* ratio (in inches) and returns a dict with keys
    ``l``, ``t``, ``r``, ``b`` (values in 1/1000 percent, i.e. 7810 ≈ 7.81%)
    suitable for hve-core's ``crop`` element attribute.  Returns ``None``
    when no meaningful crop is needed (aspect ratios already match).
    """
    from PIL import Image

    with Image.open(image_path) as img:
        iw, ih = img.size

    img_ratio = iw / ih
    tgt_ratio = target_w / target_h

    # If within 2% → no crop
    if abs(img_ratio - tgt_ratio) / max(img_ratio, tgt_ratio) < 0.02:
        return None

    if img_ratio > tgt_ratio:
        # Image wider than target → crop left/right
        visible = tgt_ratio / img_ratio
        crop_each = int(round((1 - visible) / 2 * 100_000))
        return {"l": crop_each, "r": crop_each}
    else:
        # Image taller than target → crop top/bottom
        visible = img_ratio / tgt_ratio
        crop_each = int(round((1 - visible) / 2 * 100_000))
        return {"t": crop_each, "b": crop_each}


def _image_element(
    image_name: str,
    left: float,
    top: float,
    width: float,
    height: float,
    image_path: Path | None = None,
    images_dir_name: str = "images",
) -> dict[str, Any]:
    """Build an hve-core image element dict with optional fill-crop."""
    elem: dict[str, Any] = {
        "type": "image",
        "path": f"{images_dir_name}/{image_name}",
        "left": left,
        "top": top,
        "width": width,
        "height": height,
    }
    if image_path is not None and image_path.exists():
        crop = _compute_fill_crop(image_path, width, height)
        if crop:
            elem["crop"] = crop
    return elem


def slide_to_content(
    slide: Slide,
    *,
    image_paths: list[Path] | None = None,
    label_suffix: str | None = None,
    template_path: Path | None = None,
    layouts_map: dict[str, str | int] | None = None,
    images_dir_name: str = "images",
) -> dict[str, Any]:
    """Convert a single :class:`Slide` to an hve-core ``content.yaml`` dict.

    *image_paths* are filenames already copied into the slide's ``images/``
    subfolder; only their basenames are referenced. *label_suffix*, when
    provided, is appended to the slide's title (used by draft renders to
    distinguish per-variant slides).
    """
    image_paths = image_paths or []
    layouts_map = layouts_map or {}
    layout_value = slide.layout.value
    title_text = slide.title
    if label_suffix and title_text:
        title_text = f"{title_text} {label_suffix}"
    elif label_suffix and not title_text:
        title_text = label_suffix.strip()

    roles = DEFAULT_PLACEHOLDER_ROLES.get(layout_value, PlaceholderRoles())
    placeholders: dict[int, str | list[str]] = {}
    elements: list[dict[str, Any]] = []

    # Title placeholder
    if roles.title is not None:
        v = _placeholder_value(title_text)
        if v is not None:
            placeholders[roles.title] = v

    # Body / bullets / quote text routing per layout
    if slide.layout is Layout.TITLE:
        sub = _placeholder_value(slide.subtitle)
        if sub and roles.body is not None:
            placeholders[roles.body] = sub

    elif slide.layout is Layout.SECTION_DIVIDER:
        # Section layouts have only a title placeholder; subtitle, if any,
        # is rendered as a small textbox below the title.
        sub = _placeholder_value(slide.subtitle)
        if sub:
            elements.append({
                "type": "textbox",
                "left": 0.8,
                "top": 4.6,
                "width": SLIDE_W - 1.6,
                "height": 0.7,
                "text": sub,
                "font_size": 22,
                "font_italic": True,
            })

    elif slide.layout in (Layout.BULLETS, Layout.APPENDIX_REFERENCES):
        body_lines = _bullets_to_lines(slide.bullets) or (
            [slide.body.strip()] if slide.body and slide.body.strip() else None
        )
        if body_lines and roles.body is not None:
            placeholders[roles.body] = body_lines

    elif slide.layout is Layout.BULLETS_WITH_IMAGE:
        body_lines = _bullets_to_lines(slide.bullets)
        if body_lines and roles.body is not None:
            placeholders[roles.body] = body_lines
        # Image goes into the picture-placeholder area (right half).
        if image_paths:
            dims = _picture_placeholder_dims(
                template_path,
                layouts_map.get(layout_value, DEFAULT_AZURE_LAYOUTS.get(layout_value, "")),
                roles.picture if roles.picture is not None else 10,
            ) if template_path else None
            if dims is None:
                dims = (6.67, 0.0, 6.67, SLIDE_H)
            left, top, w, h = dims
            elements.append(_image_element(
                image_paths[0].name, left, top, w, h,
                image_path=image_paths[0],
                images_dir_name=images_dir_name,
            ))

    elif slide.layout is Layout.FULL_BLEED_IMAGE:
        # Full-bleed: image covers entire canvas. Title is emitted as a
        # textbox element AFTER the image so it renders on top (z-order).
        if image_paths:
            elements.append(_image_element(
                image_paths[0].name, 0.0, 0.0, SLIDE_W, SLIDE_H,
                image_path=image_paths[0],
                images_dir_name=images_dir_name,
            ))
        # Remove title from placeholders — we add it as a textbox on top.
        placeholders.pop(roles.title, None) if roles.title is not None else None
        v = _placeholder_value(title_text)
        if v:
            # Semi-transparent dark overlay at the bottom ~25% of the slide
            # so the image breathes above while text stays readable below.
            _OVERLAY_TOP = 5.625
            _OVERLAY_H = SLIDE_H - _OVERLAY_TOP
            elements.append({
                "type": "shape",
                "shape_type": "rectangle",
                "left": 0.0,
                "top": _OVERLAY_TOP,
                "width": SLIDE_W,
                "height": _OVERLAY_H,
                "fill": {"color": "#000000", "alpha": 35},
                "line": {"width": 0},
            })
            elements.append({
                "type": "textbox",
                "left": 0.5,
                "top": _OVERLAY_TOP + 0.15,
                "width": SLIDE_W - 1.0,
                "height": _OVERLAY_H - 0.3,
                "text": v,
                "font_size": 36,
                "font_bold": True,
                "font_color": "FFFFFF",
                "vertical_alignment": "middle",
            })

    elif slide.layout in (Layout.TWO_COLUMN, Layout.COMPARISON):
        bullets = list(slide.bullets or [])
        # Even-index bullets -> left column; odd-index -> right column.
        left_col = [b for i, b in enumerate(bullets) if i % 2 == 0]
        right_col = [b for i, b in enumerate(bullets) if i % 2 == 1]
        left_lines = _bullets_to_lines(left_col)
        right_lines = _bullets_to_lines(right_col)
        if left_lines and roles.body is not None:
            placeholders[roles.body] = left_lines
        if right_lines and roles.secondary is not None:
            placeholders[roles.secondary] = right_lines

    elif slide.layout is Layout.QUOTE:
        # Opinionated layout: large italic body up top, smaller attribution
        # below. We bypass template placeholders to control vertical
        # position so longer quotes still feel balanced rather than
        # crammed into the bottom of the slide. No slide-level knobs.
        body = (slide.body or "").strip()
        attribution = _placeholder_value(slide.subtitle)

        if body:
            quoted = f"\u201c{body}\u201d"
            # Estimate body height: shrink text size for very long quotes.
            n = len(body)
            if n > 320:
                body_font = 22
            elif n > 200:
                body_font = 26
            elif n > 100:
                body_font = 30
            else:
                body_font = 36
            _Q_LEFT = 1.0
            _Q_TOP = 1.6
            _Q_W = SLIDE_W - 2.0
            _Q_H = 4.4
            elements.append({
                "type": "textbox",
                "left": _Q_LEFT,
                "top": _Q_TOP,
                "width": _Q_W,
                "height": _Q_H,
                "text": quoted,
                "font_size": body_font,
                "font_italic": True,
                "vertical_alignment": "middle",
                "alignment": "center",
            })
            if attribution:
                elements.append({
                    "type": "textbox",
                    "left": _Q_LEFT,
                    "top": _Q_TOP + _Q_H + 0.25,
                    "width": _Q_W,
                    "height": 0.6,
                    "text": f"\u2014 {attribution}",
                    "font_size": 16,
                    "font_italic": False,
                    "alignment": "center",
                })
        elif attribution:
            elements.append({
                "type": "textbox",
                "left": 1.0,
                "top": 5.5,
                "width": SLIDE_W - 2.0,
                "height": 0.6,
                "text": f"\u2014 {attribution}",
                "font_size": 16,
                "alignment": "center",
            })

    elif slide.layout is Layout.COVER:
        # Hero / title-cover: full-bleed background image, semi-transparent
        # dark rectangle on the LEFT HALF, large title and smaller subtitle
        # inside the panel. Reusable, opinionated, no per-slide knobs.
        if image_paths:
            elements.append(_image_element(
                image_paths[0].name, 0.0, 0.0, SLIDE_W, SLIDE_H,
                image_path=image_paths[0],
                images_dir_name=images_dir_name,
            ))
        _PANEL_LEFT = 0.0
        _PANEL_TOP = 0.0
        _PANEL_W = SLIDE_W * 0.5
        _PANEL_H = SLIDE_H
        elements.append({
            "type": "shape",
            "shape_type": "rectangle",
            "left": _PANEL_LEFT,
            "top": _PANEL_TOP,
            "width": _PANEL_W,
            "height": _PANEL_H,
            "fill": {"color": "#000000", "alpha": 55},
            "line": {"width": 0},
        })
        v = _placeholder_value(title_text)
        if v:
            # Long title -> shrink. The placeholder is intentionally not
            # used; layout/template title-style anchoring is unreliable
            # across templates so we own positioning here.
            n = len(v)
            if n > 80:
                title_font = 28
            elif n > 40:
                title_font = 36
            else:
                title_font = 44
            _T_LEFT = 0.6
            _T_W = _PANEL_W - 1.0
            sub = _placeholder_value(slide.subtitle)
            # Reserve room for subtitle if present.
            if sub:
                _T_TOP = 1.8
                _T_H = 3.2
            else:
                _T_TOP = 2.4
                _T_H = 2.7
            elements.append({
                "type": "textbox",
                "left": _T_LEFT,
                "top": _T_TOP,
                "width": _T_W,
                "height": _T_H,
                "text": v,
                "font_size": title_font,
                "font_bold": True,
                "font_color": "FFFFFF",
                "vertical_alignment": "bottom",
            })
            if sub:
                elements.append({
                    "type": "textbox",
                    "left": _T_LEFT,
                    "top": _T_TOP + _T_H + 0.1,
                    "width": _T_W,
                    "height": 1.4,
                    "text": sub,
                    "font_size": 18,
                    "font_color": "FFFFFF",
                    "vertical_alignment": "top",
                })

    elif slide.layout is Layout.IMAGE_GRID:
        # "Three Filmstrip Photos" has landscape PICTURE placeholders.
        _GRID_FALLBACK = [
            (0.27, 2.72, 4.40, 2.50),   # idx 13
            (4.47, 2.72, 4.40, 2.50),   # idx 14
            (8.67, 2.72, 4.40, 2.50),   # idx 15
        ]
        grid_pic_indices = [13, 14, 15]
        grid_layout_name = layouts_map.get(layout_value, DEFAULT_AZURE_LAYOUTS.get(layout_value, ""))
        grid_dims: list[tuple[float, float, float, float]] = []
        for idx in grid_pic_indices:
            d = _picture_placeholder_dims(template_path, grid_layout_name, idx) if template_path else None
            grid_dims.append(d if d else _GRID_FALLBACK[len(grid_dims)])

        chosen = image_paths[:3]
        for i, path in enumerate(chosen):
            left, top, w, h = grid_dims[i]
            elements.append(_image_element(
                path.name, left, top, w, h,
                image_path=path,
                images_dir_name=images_dir_name,
            ))

    elif slide.layout is Layout.IMAGE_SINGLE:
        # Centered landscape image below the title. Uses "Title Only" layout
        # so the title placeholder is filled normally and the image is placed
        # as a large centered element occupying most of the slide area.
        if image_paths:
            # Centre a large landscape image under the title area
            _IMG_TOP = 1.8
            _IMG_H = SLIDE_H - _IMG_TOP - 0.5
            _IMG_W = _IMG_H * 1.5  # 3:2 landscape ratio
            if _IMG_W > SLIDE_W - 1.0:
                _IMG_W = SLIDE_W - 1.0
                _IMG_H = _IMG_W / 1.5
            _IMG_LEFT = (SLIDE_W - _IMG_W) / 2
            elements.append(_image_element(
                image_paths[0].name, _IMG_LEFT, _IMG_TOP, _IMG_W, _IMG_H,
                image_path=image_paths[0],
                images_dir_name=images_dir_name,
            ))

    elif slide.layout is Layout.IMAGE_DUO:
        # Two images side by side, centered and symmetrical. Uses "Two
        # picture content" layout. We compute positions explicitly for
        # proper symmetry rather than relying on template placeholders.
        _DUO_TOP = 2.0
        _DUO_H = 4.0
        _DUO_W = _DUO_H  # square images for duo
        _DUO_GAP = 0.6
        _DUO_TOTAL = _DUO_W * 2 + _DUO_GAP
        _DUO_LEFT1 = (SLIDE_W - _DUO_TOTAL) / 2
        _DUO_LEFT2 = _DUO_LEFT1 + _DUO_W + _DUO_GAP

        duo_dims = [
            (_DUO_LEFT1, _DUO_TOP, _DUO_W, _DUO_H),
            (_DUO_LEFT2, _DUO_TOP, _DUO_W, _DUO_H),
        ]

        chosen = image_paths[:2]
        for i, path in enumerate(chosen):
            left, top, w, h = duo_dims[i]
            elements.append(_image_element(
                path.name, left, top, w, h,
                image_path=path,
                images_dir_name=images_dir_name,
            ))

    # Append the slide's own opt-in `extra_elements` passthrough. These are
    # spliced verbatim — agents/users own validity. Drawn last so they sit
    # on top of adapter-generated content.
    if slide.extra_elements:
        elements.extend(slide.extra_elements)

    content: dict[str, Any] = {
        "slide": 1,  # caller overrides
        "layout": layout_value,
    }
    if title_text:
        content["title"] = title_text
    if placeholders:
        # YAML emits int keys as integers; hve-core re-coerces via `int(idx_str)`.
        content["placeholders"] = {int(k): v for k, v in placeholders.items()}
    if elements:
        content["elements"] = elements
    if slide.notes:
        # hve-core's build_deck.py reads `speaker_notes`.
        content["speaker_notes"] = slide.notes.strip()
    return content


# ---------------------------------------------------------------------------
# Workspace materialization
# ---------------------------------------------------------------------------


def _build_style(
    pres: Presentation, *, normalized_template_path: Path | None
) -> dict[str, Any]:
    style: dict[str, Any] = {
        "dimensions": {
            "width_inches": SLIDE_W,
            "height_inches": SLIDE_H,
            "format": "16:9",
        },
    }
    if normalized_template_path is not None:
        style["template"] = {
            "path": str(normalized_template_path),
            "preserve_dimensions": True,
        }
    layouts_map: dict[str, str | int] = dict(DEFAULT_AZURE_LAYOUTS)
    layouts_map.update(pres.theme.layouts or {})
    style["layouts"] = layouts_map
    if pres.theme.metadata:
        style["metadata"] = dict(pres.theme.metadata)
    if pres.theme.defaults:
        style["defaults"] = dict(pres.theme.defaults)
    if pres.theme.fonts:
        # Surface fonts under metadata so they end up in the produced PPTX
        # properties (informational; hve-core does not consume `fonts`).
        style.setdefault("metadata", {}).setdefault(
            "subject", style["metadata"].get("subject", "")
        )
    return style


@dataclass
class WorkspacePaths:
    """Result of :func:`materialize_workspace`."""

    workdir: Path
    content_dir: Path
    style_path: Path
    template_path: Path | None
    slide_dirs: list[Path]


def materialize_workspace(
    pres: Presentation,
    *,
    workdir: Path,
    mode: str,
) -> WorkspacePaths:
    """Materialize the full hve-core content tree under *workdir*.

    *mode* must be one of ``"draft"`` or ``"final"``. ``draft`` emits one
    slide per image variant for each image-bearing slide; ``final`` emits
    one slide per spec slide, using the user's current selection.

    The returned :class:`WorkspacePaths` is everything `build_deck.py`
    needs (``content_dir``, ``style_path``, optional ``template_path``).
    """
    if mode not in ("draft", "final"):
        raise ValueError(f"mode must be 'draft' or 'final', got {mode!r}")
    workdir = Path(workdir)
    if workdir.exists():
        shutil.rmtree(workdir)
    workdir.mkdir(parents=True)
    content_dir = workdir / "content"
    global_dir = content_dir / "global"
    global_dir.mkdir(parents=True)

    template_path: Path | None = None
    if pres.theme.template:
        from .template_utils import normalize_template_to_pptx, override_layout_backgrounds

        normalized = workdir / "template.pptx"
        template_path = normalize_template_to_pptx(pres.theme.template, normalized)
        if pres.theme.layout_backgrounds:
            override_layout_backgrounds(template_path, pres.theme.layout_backgrounds)

    style = _build_style(pres, normalized_template_path=template_path)
    style_path = global_dir / "style.yaml"
    style_path.write_text(yaml.safe_dump(style, sort_keys=False), encoding="utf-8")

    slide_dirs: list[Path] = []
    slide_no = 0

    for slide in pres.slides:
        if mode == "draft" and slide.image_ref:
            variants = _list_variants(pres, slide)
            if not variants:
                slide_no += 1
                _emit_slide(
                    slide,
                    slide_no=slide_no,
                    content_dir=content_dir,
                    image_files=[],
                    label_suffix="(no images yet)",
                    template_path=template_path,
                    layouts_map=pres.theme.layouts,
                    slide_dirs=slide_dirs,
                )
                continue

            if slide.layout is Layout.IMAGE_GRID:
                # IMAGE_GRID in draft: one slide per model, each with up
                # to 3 images so the user can evaluate the grid effect.
                by_model: dict[str, list[Path]] = {}
                for mdl, _v, _i, p in variants:
                    by_model.setdefault(mdl, []).append(p)
                models = list(by_model.keys())
                for mi, mdl in enumerate(models, start=1):
                    slide_no += 1
                    suffix = f"\u2014 model {mi}/{len(models)} ({mdl})"
                    _emit_slide(
                        slide,
                        slide_no=slide_no,
                        content_dir=content_dir,
                        image_files=by_model[mdl][:3],
                        label_suffix=suffix,
                        template_path=template_path,
                        layouts_map=pres.theme.layouts,
                        slide_dirs=slide_dirs,
                    )
            elif slide.layout is Layout.IMAGE_DUO:
                # IMAGE_DUO in draft: one slide per model, each with up
                # to 2 images so the user can compare pairs.
                by_model_duo: dict[str, list[Path]] = {}
                for mdl, _v, _i, p in variants:
                    by_model_duo.setdefault(mdl, []).append(p)
                models_duo = list(by_model_duo.keys())
                for mi, mdl in enumerate(models_duo, start=1):
                    slide_no += 1
                    suffix = f"\u2014 model {mi}/{len(models_duo)} ({mdl})"
                    _emit_slide(
                        slide,
                        slide_no=slide_no,
                        content_dir=content_dir,
                        image_files=by_model_duo[mdl][:2],
                        label_suffix=suffix,
                        template_path=template_path,
                        layouts_map=pres.theme.layouts,
                        slide_dirs=slide_dirs,
                    )
            else:
                for i, (model, var, inst, src_path) in enumerate(variants, start=1):
                    slide_no += 1
                    suffix = (
                        f"\u2014 variant {i}/{len(variants)} "
                        f"({model} v{var:02d} i{inst:02d})"
                    )
                    _emit_slide(
                        slide,
                        slide_no=slide_no,
                        content_dir=content_dir,
                        image_files=[src_path],
                        label_suffix=suffix,
                        template_path=template_path,
                        layouts_map=pres.theme.layouts,
                        slide_dirs=slide_dirs,
                    )
            continue

        # Final mode (or draft for textual slides): one slide.
        image_files: list[Path] = []
        if slide.layout is Layout.IMAGE_GRID and slide.image_ref:
            # IMAGE_GRID always needs multiple images (up to 3).
            # In final mode, prefer variants from the selected model.
            sel = pres.selections.get(slide.slide_id) if mode == "final" else None
            all_variants = _list_variants(pres, slide)
            if sel:
                preferred = [p for mdl, *_, p in all_variants if mdl == sel.model.lower()]
                others = [p for mdl, *_, p in all_variants if mdl != sel.model.lower()]
                image_files = (preferred + others)[:3]
            else:
                image_files = [p for *_, p in all_variants][:3]
        elif slide.layout is Layout.IMAGE_DUO and slide.image_ref:
            # IMAGE_DUO needs 2 images. Prefer variants from the selected model.
            sel = pres.selections.get(slide.slide_id) if mode == "final" else None
            all_variants = _list_variants(pres, slide)
            if sel:
                preferred = [p for mdl, *_, p in all_variants if mdl == sel.model.lower()]
                others = [p for mdl, *_, p in all_variants if mdl != sel.model.lower()]
                image_files = (preferred + others)[:2]
            else:
                image_files = [p for *_, p in all_variants][:2]
        elif mode == "final" and slide.image_ref:
            sel_path = _resolve_selected_image(pres, slide)
            if sel_path and sel_path.exists():
                image_files = [sel_path]

        slide_no += 1
        _emit_slide(
            slide,
            slide_no=slide_no,
            content_dir=content_dir,
            image_files=image_files,
            label_suffix=None,
            template_path=template_path,
            layouts_map=pres.theme.layouts,
            slide_dirs=slide_dirs,
        )

    return WorkspacePaths(
        workdir=workdir,
        content_dir=content_dir,
        style_path=style_path,
        template_path=template_path,
        slide_dirs=slide_dirs,
    )


def _emit_slide(
    slide: Slide,
    *,
    slide_no: int,
    content_dir: Path,
    image_files: list[Path],
    label_suffix: str | None,
    template_path: Path | None,
    layouts_map: dict[str, str | int],
    slide_dirs: list[Path],
) -> None:
    slide_dir = content_dir / f"slide-{slide_no:03d}"
    images_dir = slide_dir / "images"
    slide_dir.mkdir(parents=True)
    copied: list[Path] = []
    if image_files:
        images_dir.mkdir(parents=True)
        seen_names: set[str] = set()
        for src in image_files:
            name = src.name
            if name in seen_names:
                # Disambiguate: prefix with parent dir name (model name)
                name = f"{src.parent.name}_{name}"
            seen_names.add(name)
            dst = images_dir / name
            shutil.copyfile(src, dst)
            copied.append(dst)

    content = slide_to_content(
        slide,
        image_paths=copied,
        label_suffix=label_suffix,
        template_path=template_path,
        layouts_map=layouts_map,
    )
    content["slide"] = slide_no
    (slide_dir / "content.yaml").write_text(
        yaml.safe_dump(content, sort_keys=False), encoding="utf-8"
    )
    slide_dirs.append(slide_dir)
