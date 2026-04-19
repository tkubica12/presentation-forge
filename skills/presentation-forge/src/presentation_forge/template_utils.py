"""PowerPoint template helpers.

`python-pptx` (and therefore the vendored hve-core skill) cannot open
`.potx` files directly — it expects the regular presentation content type
(`application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml`)
rather than the template content type
(`application/vnd.openxmlformats-officedocument.presentationml.template.main+xml`).

This module provides `normalize_template_to_pptx`, which copies the template
into a cache directory and patches `[Content_Types].xml` so that
`python-pptx` accepts it. The function is a no-op for files that are already
`.pptx` (it just copies them).
"""
from __future__ import annotations

import io
import shutil
import zipfile
from pathlib import Path

POTX_CONTENT_TYPE = (
    b"application/vnd.openxmlformats-officedocument.presentationml.template.main+xml"
)
PPTX_CONTENT_TYPE = (
    b"application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"
)


def normalize_template_to_pptx(template_path: Path, dest_path: Path) -> Path:
    """Copy *template_path* to *dest_path* as a `.pptx` python-pptx can open.

    If the template is already a `.pptx`, the file is copied as-is (a fresh
    copy under `dest_path`). For `.potx`, the central `[Content_Types].xml`
    entry is rewritten so the package advertises itself as a regular
    presentation. Returns the destination path.
    """
    template_path = Path(template_path)
    dest_path = Path(dest_path)
    if not template_path.is_file():
        raise FileNotFoundError(f"template not found: {template_path}")
    dest_path.parent.mkdir(parents=True, exist_ok=True)

    suffix = template_path.suffix.lower()
    if suffix == ".pptx":
        shutil.copyfile(template_path, dest_path)
        return dest_path
    if suffix != ".potx":
        raise ValueError(
            f"unsupported template extension {suffix!r}; expected .pptx or .potx"
        )

    buf = io.BytesIO()
    with zipfile.ZipFile(template_path, "r") as zin:
        with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename == "[Content_Types].xml":
                    data = data.replace(POTX_CONTENT_TYPE, PPTX_CONTENT_TYPE)
                zout.writestr(item, data)
    dest_path.write_bytes(buf.getvalue())
    return dest_path


def override_layout_backgrounds(
    pptx_path: Path,
    overrides: dict[str, str],
) -> None:
    """Override slide-layout backgrounds in-place.

    *overrides* maps layout names (e.g. ``"Photo Slide 1"``) to hex
    colour strings (e.g. ``"D9D9D6"``).  For each matched layout this
    function:

    1. Replaces the ``<p:bg>`` element with a direct solid fill so the
       background colour is immune to ``clrMapOvr`` inversions (the Azure
       template's dark-slide layouts remap ``bg2`` → ``dk2`` = blue).
    2. Resets ``<p:clrMapOvr>`` to ``<a:masterClrMapping/>`` so that text
       and accent scheme colours follow the standard (non-inverted) mapping.
    3. Injects a full-slide opaque rectangle **behind** all placeholders to
       cover any inherited master artwork (e.g. the branded grid group).
    """
    if not overrides:
        return

    from lxml import etree

    from pptx import Presentation
    from pptx.oxml.ns import qn

    NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
    NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

    prs = Presentation(str(pptx_path))
    slide_w = prs.slide_width  # EMU
    slide_h = prs.slide_height

    for layout in prs.slide_layouts:
        hex_color = overrides.get(layout.name)
        if hex_color is None:
            continue

        layout_el = layout._element  # <p:sldLayout>
        cSld = layout_el.find(qn("p:cSld"))

        # --- 1. Replace <p:bg> with a direct solid fill ----------------
        old_bg = cSld.find(qn("p:bg"))
        if old_bg is not None:
            cSld.remove(old_bg)
        new_bg_xml = (
            f'<p:bg xmlns:p="{NS_P}" xmlns:a="{NS_A}">'
            f'  <p:bgPr>'
            f'    <a:solidFill><a:srgbClr val="{hex_color}"/></a:solidFill>'
            f'    <a:effectLst/>'
            f'  </p:bgPr>'
            f'</p:bg>'
        )
        new_bg = etree.fromstring(new_bg_xml)
        # <p:bg> must be the first child of <p:cSld>
        cSld.insert(0, new_bg)

        # --- 2. Reset clrMapOvr to standard mapping -------------------
        old_cmo = layout_el.find(qn("p:clrMapOvr"))
        if old_cmo is not None:
            layout_el.remove(old_cmo)
        new_cmo_xml = (
            f'<p:clrMapOvr xmlns:p="{NS_P}" xmlns:a="{NS_A}">'
            f'  <a:masterClrMapping/>'
            f'</p:clrMapOvr>'
        )
        layout_el.append(etree.fromstring(new_cmo_xml))

        # --- 3. Fix title text colour that inherited tx2 (blue) --------
        # The master titleStyle uses schemeClr="tx2" which, under the
        # standard (non-inverted) mapping, resolves to dk2 = blue.
        # Override the title placeholder's defRPr to a direct dark colour.
        sp_tree = cSld.find(qn("p:spTree"))
        for sp in sp_tree.iter(qn("p:sp")):
            nvPr = sp.find(f"{qn('p:nvSpPr')}/{qn('p:nvPr')}")
            if nvPr is None:
                continue
            ph_el = nvPr.find(qn("p:ph"))
            if ph_el is None:
                continue
            ph_type = ph_el.get("type", "")
            ph_idx = ph_el.get("idx", "")
            # Title placeholders: type="title" or type="ctrTitle", or idx="0"
            if ph_type in ("title", "ctrTitle") or ph_idx == "0":
                txBody = sp.find(qn("p:txBody"))
                if txBody is None:
                    continue
                for defRPr in txBody.iter(qn("a:defRPr")):
                    old_fill = defRPr.find(qn("a:solidFill"))
                    if old_fill is not None:
                        defRPr.remove(old_fill)
                    fill_xml = (
                        f'<a:solidFill xmlns:a="{NS_A}">'
                        f'  <a:srgbClr val="2A1F18"/>'
                        f'</a:solidFill>'
                    )
                    defRPr.insert(0, etree.fromstring(fill_xml))

        # --- 4. Inject full-slide covering rectangle -------------------
        max_id = max(
            (int(e.get("id", "0")) for e in sp_tree.iter(qn("p:cNvPr"))),
            default=100,
        ) + 1

        rect_xml = (
            f'<p:sp xmlns:p="{NS_P}"'
            f' xmlns:a="{NS_A}"'
            f' xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
            f'<p:nvSpPr>'
            f'  <p:cNvPr id="{max_id}" name="BgOverride"/>'
            f'  <p:cNvSpPr/><p:nvPr/>'
            f'</p:nvSpPr>'
            f'<p:spPr>'
            f'  <a:xfrm><a:off x="0" y="0"/>'
            f'  <a:ext cx="{slide_w}" cy="{slide_h}"/></a:xfrm>'
            f'  <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
            f'  <a:solidFill><a:srgbClr val="{hex_color}"/></a:solidFill>'
            f'  <a:ln w="0"><a:noFill/></a:ln>'
            f'</p:spPr>'
            f'</p:sp>'
        )
        rect_el = etree.fromstring(rect_xml)
        # Insert after grpSpPr (index 1) to be the first shape.
        sp_tree.insert(2, rect_el)

    prs.save(str(pptx_path))
