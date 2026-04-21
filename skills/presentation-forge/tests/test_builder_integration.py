"""End-to-end integration test: builder -> hve-core -> .pptx.

These tests require:
  * The user's Microsoft Azure template at .templates/...potx
  * `uv` on PATH (required by the vendored hve-core skill)
  * The skills/pptx-render/.venv to have been `uv sync`ed
"""
from __future__ import annotations

import json
import shutil
from pathlib import Path

import pytest

from presentation_forge.builder import _render_one
from presentation_forge.spec import load_presentation


def _have_uv() -> bool:
    return shutil.which("uv") is not None


pytestmark = pytest.mark.skipif(
    not _have_uv(), reason="`uv` not on PATH — required to invoke vendored hve-core"
)


def test_render_final_pptx_matches_template_layouts(
    make_presentation_dir,
    populate_image_variants,
    microsoft_template_path,
    tmp_path,
):
    """A full build should produce a final.pptx whose slide layouts match
    the Microsoft template names configured in theme.yaml, and whose
    placeholders contain the expected text."""
    folder = make_presentation_dir("deck", with_template=True)
    populate_image_variants(folder, "hero-img", variations=1)
    populate_image_variants(folder, "side-img", variations=1)
    (folder / "selections.json").write_text(
        json.dumps(
            {
                "hero": {"model": "gpt-image-1", "variation": 1, "instance": 1},
                "side": {"model": "gpt-image-1", "variation": 1, "instance": 1},
            }
        ),
        encoding="utf-8",
    )
    pres = load_presentation(folder)
    pres.build_dir.mkdir(parents=True, exist_ok=True)

    out = _render_one(pres, mode="final", output=pres.final_pptx)
    assert out.exists()

    from pptx import Presentation

    prs = Presentation(str(out))
    assert len(prs.slides) == 8, "expected one rendered slide per spec slide"

    expected_layouts = [
        "Title Slide 1",
        "Photo full bleed lower title",
        "Title and Content",
        "Photo Slide 1",
        "Quote",
        "Title Only",
        "Two picture content",
        "Section Slide 1",
    ]
    actual_layouts = [s.slide_layout.name for s in prs.slides]
    assert actual_layouts == expected_layouts

    # Slide 1: title + subtitle
    s1_texts = {ph.placeholder_format.idx: ph.text_frame.text for ph in prs.slides[0].placeholders}
    assert s1_texts.get(0) == "Hello"
    assert s1_texts.get(12) == "World"

    # Slide 3: bullets — body placeholder index 10 should be 3-line text
    s3_texts = {ph.placeholder_format.idx: ph.text_frame.text for ph in prs.slides[2].placeholders}
    assert s3_texts.get(0) == "Three things"
    body = s3_texts.get(10, "")
    assert "first thing" in body and "second thing" in body and "third thing" in body

    # Slide 5: quote — rendered as textboxes (not placeholders) for better
    # vertical balance. The full slide text frame should still contain the
    # smart-quoted body and the em-dash attribution.
    s5_text = "\n".join(
        sh.text_frame.text for sh in prs.slides[4].shapes if sh.has_text_frame
    )
    assert "\u201c" in s5_text and "\u201d" in s5_text
    assert "\u2014 " in s5_text

    # Slide 6: image-single — title
    s6_texts = {ph.placeholder_format.idx: ph.text_frame.text for ph in prs.slides[5].placeholders}
    assert s6_texts.get(0) == "Solo shot"

    # Slide 7: image-duo — title
    s7_texts = {ph.placeholder_format.idx: ph.text_frame.text for ph in prs.slides[6].placeholders}
    assert s7_texts.get(0) == "Two views"

    # Slide 8: section divider title
    s8_texts = {ph.placeholder_format.idx: ph.text_frame.text for ph in prs.slides[7].placeholders}
    assert s8_texts.get(0) == "Thanks"

    # Notes preserved on slide 1
    assert prs.slides[0].notes_slide is not None
    assert "Welcome speaker notes." in prs.slides[0].notes_slide.notes_text_frame.text

    # Single master inherited from the Microsoft template
    assert len(prs.slide_masters) == 1


def test_render_draft_fans_out_per_variant(
    make_presentation_dir,
    populate_image_variants,
    microsoft_template_path,
    tmp_path,
):
    folder = make_presentation_dir("deck", with_template=True)
    populate_image_variants(folder, "hero-img", variations=3)
    populate_image_variants(folder, "side-img", variations=2)
    pres = load_presentation(folder)
    pres.build_dir.mkdir(parents=True, exist_ok=True)

    out = _render_one(pres, mode="draft", output=pres.draft_pptx)
    assert out.exists()

    from pptx import Presentation

    prs = Presentation(str(out))
    # 1 title + 3 hero variants + 1 bullets + 2 side variants + 1 quote
    # + 3 single variants + 1 duo model-group + 1 closing = 13
    assert len(prs.slides) == 13


def test_render_final_without_template_still_works(
    make_presentation_dir, tmp_path
):
    """Smoke-test: building without a template should also succeed (default master)."""
    folder = make_presentation_dir("deck", with_template=False)
    pres = load_presentation(folder)
    pres.build_dir.mkdir(parents=True, exist_ok=True)

    out = _render_one(pres, mode="final", output=pres.final_pptx)
    assert out.exists()
    from pptx import Presentation

    prs = Presentation(str(out))
    assert len(prs.slides) == 8
