"""PPTX renderer using python-pptx.

We build slides programmatically rather than cloning template layouts by
name — this avoids fragile template assumptions and makes the renderer
work with any reasonable starting template (or no template at all). The
template, when provided, supplies slide size + theme colors via the
underlying XML; we draw all content as new shapes on blank layouts.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation as PptxPresentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from .layouts import Layout
from .slides_parser import Slide
from .spec import Presentation, Selection, Theme


def _hex_to_rgb(hex_str: str) -> RGBColor:
    s = hex_str.lstrip("#")
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


def _theme_color(theme: Theme, key: str, default: str) -> RGBColor:
    return _hex_to_rgb(theme.colors.get(key, default))


def _heading_font(theme: Theme) -> str:
    return theme.fonts.get("heading", "Calibri")


def _body_font(theme: Theme) -> str:
    return theme.fonts.get("body", "Calibri")


def _new_pptx(theme: Theme) -> PptxPresentation:
    if theme.template and theme.template.exists():
        prs = PptxPresentation(str(theme.template))
    else:
        prs = PptxPresentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    return prs


def _blank_layout(prs: PptxPresentation):
    # Layout 6 is "Blank" in default templates; fall back to 0 if not.
    try:
        return prs.slide_layouts[6]
    except IndexError:
        return prs.slide_layouts[0]


def _add_text_box(slide, *, x, y, w, h, text, font, size_pt, bold=False,
                  color: RGBColor | None = None, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    if color is not None:
        run.font.color.rgb = color
    return box


def _add_bullets_box(slide, *, x, y, w, h, bullets, font, size_pt, color):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = f"•  {bullet}"
        run.font.name = font
        run.font.size = Pt(size_pt)
        run.font.color.rgb = color
    return box


def _add_background(slide, prs, color: RGBColor):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    # Send to back
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def _add_speaker_notes(slide, notes: str | None) -> None:
    if not notes:
        return
    nf = slide.notes_slide.notes_text_frame
    nf.text = notes


# ---------------------------------------------------------------------------
# Per-layout renderers
# ---------------------------------------------------------------------------

def _render_title(slide_obj, prs, slide: Slide, theme: Theme, image_path: Path | None):
    bg = _theme_color(theme, "background", "0F172A")
    accent = _theme_color(theme, "accent", "38BDF8")
    fg = _theme_color(theme, "foreground", "F8FAFC")
    _add_background(slide_obj, prs, bg)
    margin = Inches(0.8)
    width = prs.slide_width - 2 * margin
    _add_text_box(
        slide_obj, x=margin, y=Inches(2.5), w=width, h=Inches(1.5),
        text=slide.title or "", font=_heading_font(theme), size_pt=54, bold=True,
        color=fg,
    )
    if slide.subtitle:
        _add_text_box(
            slide_obj, x=margin, y=Inches(4.2), w=width, h=Inches(0.8),
            text=slide.subtitle, font=_body_font(theme), size_pt=22,
            color=accent, italic=True,
        )


def _render_section_divider(slide_obj, prs, slide: Slide, theme: Theme, image_path):
    accent = _theme_color(theme, "accent", "38BDF8")
    fg = _theme_color(theme, "foreground", "F8FAFC")
    _add_background(slide_obj, prs, accent)
    margin = Inches(0.8)
    _add_text_box(
        slide_obj, x=margin, y=Inches(2.8), w=prs.slide_width - 2 * margin,
        h=Inches(2.0), text=slide.title or "", font=_heading_font(theme),
        size_pt=60, bold=True, color=fg,
    )
    if slide.subtitle:
        _add_text_box(
            slide_obj, x=margin, y=Inches(4.6), w=prs.slide_width - 2 * margin,
            h=Inches(0.8), text=slide.subtitle, font=_body_font(theme),
            size_pt=22, color=fg, italic=True,
        )


def _render_bullets(slide_obj, prs, slide: Slide, theme: Theme, image_path):
    bg = _theme_color(theme, "background", "FFFFFF")
    fg = _theme_color(theme, "foreground", "1E293B")
    _add_background(slide_obj, prs, bg)
    margin = Inches(0.8)
    width = prs.slide_width - 2 * margin
    _add_text_box(
        slide_obj, x=margin, y=Inches(0.6), w=width, h=Inches(1.0),
        text=slide.title or "", font=_heading_font(theme), size_pt=36, bold=True,
        color=fg,
    )
    _add_bullets_box(
        slide_obj, x=margin, y=Inches(1.8), w=width, h=Inches(5.0),
        bullets=slide.bullets, font=_body_font(theme), size_pt=20, color=fg,
    )


def _render_bullets_with_image(slide_obj, prs, slide: Slide, theme: Theme, image_path):
    bg = _theme_color(theme, "background", "FFFFFF")
    fg = _theme_color(theme, "foreground", "1E293B")
    _add_background(slide_obj, prs, bg)
    margin = Inches(0.6)
    title_h = Inches(1.0)
    _add_text_box(
        slide_obj, x=margin, y=Inches(0.5), w=prs.slide_width - 2 * margin,
        h=title_h, text=slide.title or "", font=_heading_font(theme), size_pt=32,
        bold=True, color=fg,
    )
    body_top = Inches(1.7)
    body_h = prs.slide_height - body_top - Inches(0.5)
    half = (prs.slide_width - 2 * margin - Inches(0.4)) / 2
    image_position = (slide.image_position or "right").lower()
    if image_position == "left":
        img_x, txt_x = margin, margin + half + Inches(0.4)
    else:
        txt_x, img_x = margin, margin + half + Inches(0.4)
    _add_bullets_box(
        slide_obj, x=txt_x, y=body_top, w=half, h=body_h,
        bullets=slide.bullets, font=_body_font(theme), size_pt=18, color=fg,
    )
    if image_path and image_path.exists():
        slide_obj.shapes.add_picture(
            str(image_path), img_x, body_top, width=half, height=body_h
        )


def _render_full_bleed_image(slide_obj, prs, slide: Slide, theme: Theme, image_path):
    bg = _theme_color(theme, "background", "000000")
    fg = _theme_color(theme, "foreground", "FFFFFF")
    _add_background(slide_obj, prs, bg)
    if image_path and image_path.exists():
        slide_obj.shapes.add_picture(
            str(image_path), 0, 0, width=prs.slide_width, height=prs.slide_height
        )
    if slide.title:
        # overlay box at the bottom
        overlay_h = Inches(1.4)
        overlay_y = prs.slide_height - overlay_h
        overlay = slide_obj.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, overlay_y, prs.slide_width, overlay_h
        )
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = bg
        overlay.fill.fore_color.brightness = 0
        overlay.line.fill.background()
        # Pseudo-transparency by manual alpha on the fill XML would be ideal;
        # for simplicity we use a solid dark band.
        _add_text_box(
            slide_obj, x=Inches(0.6), y=overlay_y + Inches(0.3),
            w=prs.slide_width - Inches(1.2), h=Inches(0.8),
            text=slide.title, font=_heading_font(theme), size_pt=32,
            bold=True, color=fg,
        )


def _render_two_column(slide_obj, prs, slide: Slide, theme: Theme, image_path):
    bg = _theme_color(theme, "background", "FFFFFF")
    fg = _theme_color(theme, "foreground", "1E293B")
    _add_background(slide_obj, prs, bg)
    margin = Inches(0.8)
    _add_text_box(
        slide_obj, x=margin, y=Inches(0.5), w=prs.slide_width - 2 * margin,
        h=Inches(1.0), text=slide.title or "", font=_heading_font(theme),
        size_pt=32, bold=True, color=fg,
    )
    body_top = Inches(1.8)
    body_h = prs.slide_height - body_top - Inches(0.5)
    cols = slide.bullets[:2] + [""] * max(0, 2 - len(slide.bullets))
    half = (prs.slide_width - 2 * margin - Inches(0.4)) / 2
    for i, text in enumerate(cols):
        x = margin + (half + Inches(0.4)) * i
        _add_text_box(
            slide_obj, x=x, y=body_top, w=half, h=body_h, text=text,
            font=_body_font(theme), size_pt=20, color=fg,
        )


def _render_quote(slide_obj, prs, slide: Slide, theme: Theme, image_path):
    bg = _theme_color(theme, "background", "0F172A")
    fg = _theme_color(theme, "foreground", "F8FAFC")
    accent = _theme_color(theme, "accent", "38BDF8")
    _add_background(slide_obj, prs, bg)
    margin = Inches(1.5)
    _add_text_box(
        slide_obj, x=margin, y=Inches(2.0), w=prs.slide_width - 2 * margin,
        h=Inches(3.5), text=f"\u201c{(slide.body or '').strip()}\u201d",
        font=_heading_font(theme), size_pt=36, italic=True, color=fg,
        align=PP_ALIGN.CENTER,
    )
    if slide.subtitle:
        _add_text_box(
            slide_obj, x=margin, y=Inches(5.5), w=prs.slide_width - 2 * margin,
            h=Inches(0.8), text=f"\u2014 {slide.subtitle}",
            font=_body_font(theme), size_pt=18, color=accent,
            align=PP_ALIGN.CENTER,
        )


def _render_comparison(slide_obj, prs, slide: Slide, theme: Theme, image_path):
    # same as two-column for v0.1
    _render_two_column(slide_obj, prs, slide, theme, image_path)


def _render_image_grid(slide_obj, prs, slide: Slide, theme: Theme, image_path):
    # placeholder — image grid is generated for draft.pptx by the draft builder,
    # not authored. Treat as bullets fallback.
    _render_bullets(slide_obj, prs, slide, theme, image_path)


def _render_appendix_references(slide_obj, prs, slide: Slide, theme: Theme, image_path):
    bg = _theme_color(theme, "background", "FFFFFF")
    fg = _theme_color(theme, "foreground", "1E293B")
    _add_background(slide_obj, prs, bg)
    margin = Inches(0.8)
    _add_text_box(
        slide_obj, x=margin, y=Inches(0.5), w=prs.slide_width - 2 * margin,
        h=Inches(1.0), text=slide.title or "References",
        font=_heading_font(theme), size_pt=32, bold=True, color=fg,
    )
    _add_text_box(
        slide_obj, x=margin, y=Inches(1.8), w=prs.slide_width - 2 * margin,
        h=Inches(5.0), text=slide.body or "", font=_body_font(theme),
        size_pt=14, color=fg,
    )


_RENDERERS = {
    Layout.TITLE: _render_title,
    Layout.SECTION_DIVIDER: _render_section_divider,
    Layout.BULLETS: _render_bullets,
    Layout.BULLETS_WITH_IMAGE: _render_bullets_with_image,
    Layout.FULL_BLEED_IMAGE: _render_full_bleed_image,
    Layout.TWO_COLUMN: _render_two_column,
    Layout.QUOTE: _render_quote,
    Layout.COMPARISON: _render_comparison,
    Layout.IMAGE_GRID: _render_image_grid,
    Layout.APPENDIX_REFERENCES: _render_appendix_references,
}


def _render_one(prs: PptxPresentation, slide: Slide, theme: Theme,
                image_path: Path | None, *, label_suffix: str | None = None) -> None:
    slide_obj = prs.slides.add_slide(_blank_layout(prs))
    renderer = _RENDERERS[slide.layout]
    # Optionally append a label suffix to the title for draft variants.
    rendered_slide = slide
    if label_suffix and slide.title:
        from dataclasses import replace
        rendered_slide = replace(slide, title=f"{slide.title} {label_suffix}")
    elif label_suffix and slide.layout in {Layout.FULL_BLEED_IMAGE, Layout.QUOTE}:
        from dataclasses import replace
        rendered_slide = replace(slide, title=label_suffix.strip())
    renderer(slide_obj, prs, rendered_slide, theme, image_path)
    _add_speaker_notes(slide_obj, slide.notes)


def _resolve_selected_image(pres: Presentation, slide: Slide) -> Path | None:
    if not slide.image_ref:
        return None
    sel = pres.selections.get(slide.slide_id)
    if sel is None:
        return None
    return (
        pres.images_dir / slide.image_ref / sel.model.lower()
        / sel.filename(slide.image_ref)
    )


def _list_all_variants(pres: Presentation, slide: Slide) -> list[tuple[str, int, int, Path]]:
    """Return [(model, var, inst, path), ...] for an image_ref's existing PNGs."""
    out: list[tuple[str, int, int, Path]] = []
    if not slide.image_ref:
        return out
    base = pres.images_dir / slide.image_ref
    if not base.exists():
        return out
    import re
    pat = re.compile(rf"^{re.escape(slide.image_ref)}_v(\d+)_i(\d+)\.png$")
    for model_dir in sorted(p for p in base.iterdir() if p.is_dir()):
        for png in sorted(model_dir.iterdir()):
            m = pat.match(png.name)
            if m:
                out.append((model_dir.name, int(m.group(1)), int(m.group(2)), png))
    return out


def render_final(pres: Presentation) -> Path:
    prs = _new_pptx(pres.theme)
    # python-pptx loaded templates often include a default first slide; trim.
    _strip_existing_slides(prs)
    for slide in pres.slides:
        img = _resolve_selected_image(pres, slide)
        _render_one(prs, slide, pres.theme, img)
    pres.build_dir.mkdir(parents=True, exist_ok=True)
    prs.save(str(pres.final_pptx))
    return pres.final_pptx


def render_draft(pres: Presentation) -> Path:
    prs = _new_pptx(pres.theme)
    _strip_existing_slides(prs)
    for slide in pres.slides:
        if slide.image_ref:
            variants = _list_all_variants(pres, slide)
            if not variants:
                _render_one(prs, slide, pres.theme, None, label_suffix="(no images yet)")
                continue
            for i, (model, var, inst, path) in enumerate(variants, start=1):
                suffix = f"\u2014 variant {i}/{len(variants)} ({model} v{var:02d} i{inst:02d})"
                _render_one(prs, slide, pres.theme, path, label_suffix=suffix)
        else:
            _render_one(prs, slide, pres.theme, None)
    pres.build_dir.mkdir(parents=True, exist_ok=True)
    prs.save(str(pres.draft_pptx))
    return pres.draft_pptx


def _strip_existing_slides(prs: PptxPresentation) -> None:
    """Remove any pre-existing slides from a template-loaded presentation."""
    sldIdLst = prs.slides._sldIdLst  # noqa: SLF001
    rIds = [sld.rId for sld in list(sldIdLst)]
    for sld in list(sldIdLst):
        sldIdLst.remove(sld)
    for rId in rIds:
        try:
            prs.part.drop_rel(rId)
        except KeyError:
            pass
