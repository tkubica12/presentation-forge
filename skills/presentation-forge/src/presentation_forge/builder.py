"""Build orchestration: validate spec, run image-generator, render PPTX(s)."""
from __future__ import annotations

import os
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

import yaml

from . import render_adapter
from .layouts import REQUIRED_FIELDS
from .spec import Presentation
from .state import State, hash_slide, hash_text


SKILL_DIR = Path(__file__).resolve().parents[2]  # skills/presentation-forge/
SKILLS_ROOT = SKILL_DIR.parent  # skills/
IMAGE_GENERATOR_DIR = SKILLS_ROOT / "image-generator"
PPTX_RENDER_DIR = SKILLS_ROOT / "pptx-render"
BUILD_DECK_SCRIPT = PPTX_RENDER_DIR / "scripts" / "build_deck.py"


class ValidationError(Exception):
    pass


def validate(pres: Presentation) -> list[str]:
    """Return list of validation warnings; raise on hard errors."""
    warnings: list[str] = []
    image_names = pres.image_names_in_yaml()

    for slide in pres.slides:
        # image_ref must resolve
        if slide.image_ref and slide.image_ref not in image_names:
            raise ValidationError(
                f"slide {slide.slide_id!r}: image_ref {slide.image_ref!r} "
                f"not found in images.yaml"
            )
        # required fields per layout — already checked at parse, double-check
        missing = REQUIRED_FIELDS[slide.layout] - set(slide.raw.keys())
        if missing:
            raise ValidationError(
                f"slide {slide.slide_id!r}: missing required fields {sorted(missing)}"
            )

    # Theme template, if specified, must exist
    if pres.theme.template and not pres.theme.template.exists():
        raise ValidationError(
            f"theme.yaml template path does not resolve: {pres.theme.template}"
        )

    # Selections referencing files that don't exist (yet) — warn only.
    for slide in pres.slides:
        sel = pres.selections.get(slide.slide_id)
        if sel is None or not slide.image_ref:
            continue
        path = (
            pres.images_dir / slide.image_ref / sel.model.lower()
            / sel.filename(slide.image_ref)
        )
        if not path.exists():
            warnings.append(
                f"selection for {slide.slide_id!r} points at missing file: {path}"
            )

    return warnings


def _find_uv() -> str:
    uv = shutil.which("uv")
    if not uv:
        raise FileNotFoundError("`uv` not found in PATH; install uv to run sub-skills")
    return uv


def _fix_fullbleed_zorder(pptx_path: Path) -> None:
    """Move full-canvas picture shapes behind all other shapes.

    hve-core adds image *elements* after placeholders, so they sit on
    top of the title text.  This post-pass reorders full-bleed images to
    be the first shapes in the shape tree, restoring the expected z-order.
    """
    from pptx import Presentation as PptxPresentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = PptxPresentation(str(pptx_path))
    canvas_w = prs.slide_width
    canvas_h = prs.slide_height
    threshold = 0.85

    changed = False
    for slide in prs.slides:
        sp_tree = slide.shapes._spTree
        for shape in list(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                w_ratio = shape.width / canvas_w if canvas_w else 0
                h_ratio = shape.height / canvas_h if canvas_h else 0
                if w_ratio > threshold and h_ratio > threshold:
                    elem = shape._element
                    sp_tree.remove(elem)
                    # Insert after nvGrpSpPr (idx 0) and grpSpPr (idx 1)
                    sp_tree.insert(2, elem)
                    changed = True

    if changed:
        prs.save(str(pptx_path))


def run_images(
    pres: Presentation,
    *,
    parallelism: int | None = None,
    only: list[str] | None = None,
) -> None:
    """Shell out to the sibling image-generator skill.

    *only*: optional list of image_ref names to restrict generation to. When
    provided, a temporary filtered images.yaml is written and passed to the
    image-generator so other entries are skipped entirely.
    """
    if not IMAGE_GENERATOR_DIR.exists():
        raise FileNotFoundError(
            f"sibling image-generator skill not found at {IMAGE_GENERATOR_DIR}. "
            f"Install it with: gh skill install tkubica12/presentation-forge image-generator"
        )

    yaml_path = pres.images_yaml_path
    tmp_yaml: Path | None = None
    if only:
        only_set = {o.strip() for o in only if o and o.strip()}
        all_imgs = pres.images_yaml_data.get("images") or []
        filtered = [img for img in all_imgs if img.get("name") in only_set]
        unknown = only_set - {img.get("name") for img in all_imgs}
        if unknown:
            raise ValueError(
                f"--only references unknown image_ref(s): {sorted(unknown)}"
            )
        if not filtered:
            print("nothing to do (no matching image_refs)", flush=True)
            return
        data = dict(pres.images_yaml_data)
        data["images"] = filtered
        tmp = tempfile.NamedTemporaryFile(
            mode="w", suffix=".yaml", delete=False, encoding="utf-8"
        )
        tmp.write(yaml.safe_dump(data, sort_keys=False, allow_unicode=True))
        tmp.close()
        tmp_yaml = Path(tmp.name)
        yaml_path = tmp_yaml
        print(
            f"image-generator scope: only {sorted(only_set)} "
            f"({len(filtered)} of {len(all_imgs)} entries)",
            flush=True,
        )

    uv = _find_uv()
    cmd = [
        uv, "--directory", str(IMAGE_GENERATOR_DIR), "run", "generate-images",
        str(yaml_path),
        "--output-dir", str(pres.images_dir),
    ]
    if parallelism is not None:
        cmd.extend(["--parallelism", str(parallelism)])
    print(f"$ {' '.join(cmd)}", flush=True)
    try:
        rc = subprocess.call(cmd, env=_subprocess_env())
    finally:
        if tmp_yaml is not None:
            try:
                tmp_yaml.unlink()
            except OSError:
                pass
    if rc != 0:
        raise RuntimeError(f"image-generator exited with code {rc}")


def _subprocess_env() -> dict[str, str]:
    """Subprocess env with UTF-8 stdio so non-ASCII titles don't crash on Windows."""
    env = os.environ.copy()
    env.setdefault("PYTHONIOENCODING", "utf-8")
    env.setdefault("PYTHONUTF8", "1")
    return env


def _is_locked(path: Path) -> bool:
    """Best-effort detection of a Windows-locked output file.

    On Windows, opening a file that PowerPoint holds with append-mode
    will raise PermissionError. Returns True if the file is locked.
    """
    if not path.exists():
        return False
    try:
        with open(path, "a"):
            pass
        return False
    except PermissionError:
        return True
    except OSError:
        return False


def _alternate_output(path: Path) -> Path:
    """Return ``<stem>-updated[.<n>].<suffix>`` next to *path*."""
    stem, suffix, parent = path.stem, path.suffix, path.parent
    candidate = parent / f"{stem}-updated{suffix}"
    n = 2
    while candidate.exists() and _is_locked(candidate):
        candidate = parent / f"{stem}-updated-{n}{suffix}"
        n += 1
    return candidate


def _render_one(pres: Presentation, *, mode: str, output: Path) -> Path:
    """Materialize hve-core workspace and shell out to its build_deck.py."""
    if not BUILD_DECK_SCRIPT.exists():
        raise FileNotFoundError(
            f"sibling pptx-render skill not found at {PPTX_RENDER_DIR}. "
            f"It should be vendored at skills/pptx-render/."
        )
    workdir = pres.build_dir / f"_hve_{mode}"
    paths = render_adapter.materialize_workspace(pres, workdir=workdir, mode=mode)

    # If the desired output is open in PowerPoint (Windows lock), divert
    # to a sibling filename so the user doesn't lose the rebuild.
    actual_output = output
    if _is_locked(output):
        actual_output = _alternate_output(output)
        print(
            f"⚠  {output.name} is locked (open in PowerPoint?). "
            f"Writing to {actual_output.name} instead.",
            flush=True,
        )

    actual_output.parent.mkdir(parents=True, exist_ok=True)
    uv = _find_uv()
    cmd = [
        uv, "--directory", str(PPTX_RENDER_DIR), "run",
        "python", str(BUILD_DECK_SCRIPT),
        "--content-dir", str(paths.content_dir),
        "--style", str(paths.style_path),
        "--output", str(actual_output),
    ]
    if paths.template_path is not None:
        cmd.extend(["--template", str(paths.template_path)])
    print(f"$ {' '.join(cmd)}", flush=True)
    rc = subprocess.call(cmd, env=_subprocess_env())
    if rc != 0:
        raise RuntimeError(f"build_deck.py exited with code {rc}")
    if not actual_output.exists():
        raise RuntimeError(f"build_deck.py did not produce expected output: {actual_output}")
    _fix_fullbleed_zorder(actual_output)
    return actual_output


def build(pres: Presentation, *, draft: bool = True, final: bool = True,
          run_image_gen: bool = True, only: list[str] | None = None) -> dict[str, Path]:
    """Run the full build. Returns paths of generated artifacts.

    *only*: optional list of image_ref names; when provided, image generation
    is restricted to those refs (other refs are NOT regenerated and existing
    PNGs are reused). Has no effect on the rendered PPTX (which always
    includes all slides).
    """
    if run_image_gen:
        run_images(pres, only=only)
    pres.build_dir.mkdir(parents=True, exist_ok=True)
    state = State.load(pres.state_path)

    out: dict[str, Path] = {}
    if draft:
        out["draft"] = _render_one(pres, mode="draft", output=pres.draft_pptx)
        for slide in pres.slides:
            state.record_slide(slide.slide_id, hash_slide(slide), kind="draft")
    if final:
        out["final"] = _render_one(pres, mode="final", output=pres.final_pptx)
        for slide in pres.slides:
            state.record_slide(slide.slide_id, hash_slide(slide), kind="final")

    state.images.setdefault("yaml_hash", hash_text(
        pres.images_yaml_path.read_text(encoding="utf-8")
    ))
    state.save(pres.state_path)
    return out


def regenerate_image(pres: Presentation, image_ref: str) -> None:
    """Wipe the cache for *image_ref* and regenerate from scratch.

    Removes ``build/images/<image_ref>/`` (PNGs + prompts.json) so the
    image-generator regenerates everything (including prompts) instead of
    reusing stale outputs.
    """
    if image_ref not in pres.image_names_in_yaml():
        raise ValueError(f"unknown image_ref {image_ref!r}; not in images.yaml")
    target = pres.images_dir / image_ref
    if target.exists():
        shutil.rmtree(target)
        print(f"removed {target}", flush=True)
    run_images(pres, only=[image_ref])


def images_status(pres: Presentation) -> list[dict]:
    """Per-image_ref status: expected variants × instances vs PNGs found."""
    yaml_data = pres.images_yaml_data
    var_count = int(yaml_data.get("variations_count")
                    or yaml_data.get("variations_per_image") or 4)
    inst_count = int(yaml_data.get("instances_per_prompt")
                     or yaml_data.get("instances_per_variation") or 1)
    rows: list[dict] = []
    for entry in (yaml_data.get("images") or []):
        ref = entry.get("name")
        if not ref:
            continue
        base = pres.images_dir / ref
        models: list[str] = []
        png_count = 0
        if base.exists():
            for model_dir in sorted(p for p in base.iterdir() if p.is_dir()):
                models.append(model_dir.name)
                png_count += sum(
                    1 for f in model_dir.iterdir() if f.suffix.lower() == ".png"
                )
        expected_per_model = var_count * inst_count
        rows.append({
            "image_ref": ref,
            "models": models,
            "pngs_found": png_count,
            "expected_per_model": expected_per_model,
            "complete": bool(models)
                and all(
                    sum(
                        1 for f in (base / m).iterdir() if f.suffix.lower() == ".png"
                    ) >= expected_per_model
                    for m in models
                ),
        })
    return rows


def status(pres: Presentation) -> list[dict]:
    """Per-slide summary used by `forge status`."""
    state = State.load(pres.state_path)
    out: list[dict] = []
    for slide in pres.slides:
        h = hash_slide(slide)
        prev = state.slides.get(slide.slide_id, {})
        sel = pres.selections.get(slide.slide_id)
        out.append({
            "slide_id": slide.slide_id,
            "layout": slide.layout.value,
            "image_ref": slide.image_ref,
            "selection": (
                f"{sel.model} v{sel.variation:02d} i{sel.instance:02d}"
                if sel else None
            ),
            "changed_since_last_build": prev.get("content_hash") != h if prev else True,
            "last_built": prev.get("last_built_final") or prev.get("last_built_draft"),
        })
    return out
